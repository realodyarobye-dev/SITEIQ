"""PowerPoint deck.

Fifteen designed slides for showing a partner, a landlord, an investor or your
family. Dark editorial styling, one idea per slide, no paragraph dumps.
"""
import io
import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

log = logging.getLogger("siteiq.pptx")

BG = RGBColor(0x0B, 0x14, 0x20)
PANEL = RGBColor(0x14, 0x22, 0x33)
PANEL2 = RGBColor(0x1B, 0x2C, 0x40)
WHITE = RGBColor(0xF4, 0xF8, 0xFB)
MUTED = RGBColor(0x93, 0xA6, 0xBA)
MINT = RGBColor(0x5E, 0xE0, 0xA0)

VERDICT_RGB = {
    "TAKE IT": RGBColor(0x18, 0xA8, 0x70), "STRONG": RGBColor(0x4C, 0xAE, 0x5A),
    "NEGOTIATE": RGBColor(0xE0, 0xA0, 0x2A), "MAYBE": RGBColor(0xDD, 0x82, 0x33),
    "PASS": RGBColor(0xD1, 0x4B, 0x4B),
}
SEV_RGB = {"CRITICAL": RGBColor(0xD1, 0x4B, 0x4B), "HIGH": RGBColor(0xE0, 0x74, 0x2A),
           "MEDIUM": RGBColor(0xE0, 0xA0, 0x2A), "LOW": RGBColor(0x7D, 0x8C, 0xA0)}
THREAT_RGB = SEV_RGB | {"VERY HIGH": RGBColor(0xD1, 0x4B, 0x4B)}

W, H = Inches(13.333), Inches(7.5)


def build(report, photos=None):
    photos = photos or {}
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    _cover(prs, report)
    _verdict(prs, report)
    _overview(prs, report)
    _block(prs, report)
    _landscape(prs, report)
    _competitors(prs, report)
    _generators(prs, report)
    _customers(prs, report)
    _dayparts(prs, report)
    _demographics(prs, report)
    _sales(prs, report)
    _economics(prs, report)
    _opportunity(prs, report)
    _risks(prs, report)
    _final(prs, report, photos)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ------------------------------------------------------------------ helpers
def _slide(prs, title=None, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, BG)
    _rect(s, 0, H - Inches(0.06), W, Inches(0.06), MINT)
    if title:
        _text(s, 0.55, 0.34, 12.2, 0.62, title, 30, WHITE, bold=True)
    if sub:
        _text(s, 0.58, 1.02, 12.2, 0.42, sub, 13, MUTED)
    return s


def _rect(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(x), _emu(y), _emu(w), _emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _emu(v):
    return v if isinstance(v, Emu) else Inches(v)


def _text(slide, x, y, w, h, text, size, colour, bold=False, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = "Calibri"
    return box


def _bullets(slide, x, y, w, h, items, size=15, colour=WHITE, spacing=8):
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = str(head)
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = colour
            r2 = p.add_run(); r2.text = "  " + str(tail)
            r2.font.size = Pt(size - 2); r2.font.color.rgb = MUTED
        else:
            r = p.add_run(); r.text = str(item)
            r.font.size = Pt(size); r.font.color.rgb = colour
    return box


def _kpi(slide, x, y, w, h, label, value, colour=WHITE, sub=None):
    _rect(slide, x, y, w, h, PANEL)
    _text(slide, x + 0.22, y + 0.16, w - 0.4, 0.3, label.upper(), 9.5, MUTED, bold=True)
    _text(slide, x + 0.2, y + 0.46, w - 0.35, 0.6, value, 25, colour, bold=True)
    if sub:
        _text(slide, x + 0.22, y + h - 0.42, w - 0.4, 0.3, sub, 10, MUTED)


def _bar(slide, x, y, w, h, pct, colour=MINT):
    _rect(slide, x, y, w, h, PANEL2)
    filled = max(0.02, w * min(100, max(0, pct)) / 100.0)
    _rect(slide, x, y, filled, h, colour)


def _table(slide, x, y, w, headers, rows, col_widths=None, size=11):
    cols = len(headers)
    rowcount = len(rows) + 1
    shape = slide.shapes.add_table(rowcount, cols, _emu(x), _emu(y), _emu(w),
                                   Inches(0.34 * rowcount))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Emu(int(_emu(w) * cw / total))
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(head)
        _cell_style(cell, 10, MUTED, PANEL2, bold=True)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            colour = WHITE
            if isinstance(val, tuple):
                val, colour = val
            cell.text = str(val)
            _cell_style(cell, size, colour, PANEL if r % 2 else BG)
    return table


def _cell_style(cell, size, colour, fill, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.09)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = colour
            r.font.bold = bold
            r.font.name = "Calibri"


def _note(slide, text):
    _text(slide, 0.55, 6.86, 12.2, 0.35, text, 9, MUTED, italic=True)


# ------------------------------------------------------------------- slides
def _cover(prs, r):
    s = _slide(prs)
    _rect(s, 0, 0, W, H, BG)
    _rect(s, 0, 0, Inches(0.13), H, VERDICT_RGB.get(r["verdict"], MINT))
    _rect(s, 0, H - Inches(0.06), W, Inches(0.06), MINT)
    _text(s, 0.75, 1.5, 8, 0.5, "SITEIQ", 20, MINT, bold=True)
    _text(s, 0.72, 2.0, 11.5, 1.4, r["address"], 42, WHITE, bold=True)
    _text(s, 0.75, 3.4, 11.5, 0.5, r["resolved_address"], 14, MUTED)
    _text(s, 0.75, 4.05, 11.5, 0.4, r["concept_label"].upper(), 13, MINT, bold=True)

    _kpi(s, 0.72, 4.7, 3.3, 1.35, "Verdict", r["verdict"], VERDICT_RGB.get(r["verdict"], WHITE))
    _kpi(s, 4.22, 4.7, 2.6, 1.35, "Overall score", f"{r['score']}", WHITE, "out of 100")
    _kpi(s, 7.02, 4.7, 2.6, 1.35, "Data confidence", f"{r['confidence']['score']}", WHITE, "out of 100")
    _kpi(s, 9.82, 4.7, 2.7, 1.35, "Realistic sales",
         f"${r['sales']['scenarios']['Realistic']['daily']:,}", MINT, "per day, modeled")
    _note(s, "All sales, profit and traffic figures are modeled estimates from public data, not "
             "measured results.")


def _verdict(prs, r):
    s = _slide(prs, "Executive verdict", r["resolved_address"])
    colour = VERDICT_RGB.get(r["verdict"], MINT)
    _rect(s, 0.55, 1.62, 12.2, 0.95, colour)
    _text(s, 0.8, 1.78, 11.7, 0.65, f"{r['verdict']}  -  {r['five_hundred_k']['headline']}",
          20, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    _text(s, 0.58, 2.78, 12.2, 1.0, r["verdict_reason"], 14, WHITE)

    sc = r["sales"]["scenarios"]
    pnl = r["pnl"]["by_scenario"]
    _kpi(s, 0.55, 4.0, 2.9, 1.25, "Conservative", f"${sc['Conservative']['daily']:,}", MUTED, "per day")
    _kpi(s, 3.65, 4.0, 2.9, 1.25, "Realistic", f"${sc['Realistic']['daily']:,}", MINT, "per day")
    _kpi(s, 6.75, 4.0, 2.9, 1.25, "Strong operator", f"${sc['Strong Operator']['daily']:,}", WHITE, "per day")
    _kpi(s, 9.85, 4.0, 2.9, 1.25, "Elite operator", f"${sc['Elite Operator']['daily']:,}", WHITE, "per day")

    profit = pnl["Realistic"]["profit"]
    _kpi(s, 0.55, 5.42, 4.0, 1.25, "Monthly operating profit (realistic)",
         f"${profit:,}", MINT if profit > 0 else SEV_RGB["CRITICAL"],
         f"{pnl['Realistic']['margin_pct']}% margin")
    _kpi(s, 4.75, 5.42, 4.0, 1.25, "Rent",
         f"${r['rent']['monthly_rent']:,}", WHITE,
         f"{r['rent_assessment']['rent_pct']}% of sales - {r['rent_assessment']['band']}"
         + (" (ESTIMATED)" if r["rent"]["estimated"] else ""))
    _kpi(s, 8.95, 5.42, 3.8, 1.25, "Competitors within 0.5 mi",
         f"{r['competition']['direct_count']} direct", WHITE,
         r["competition"]["saturation_label"])
    _note(s, "MODELED. Estimates only.")


def _overview(prs, r):
    s = _slide(prs, "Location overview", "What is around this address")
    rows = []
    for x in r["radii"]:
        top = ", ".join(f"{k} ({v})" for k, v in list(x["by_category"].items())[:3])
        rows.append([x["label"], f"{x['total']}", f"{x['food']}", top or "-"])
    _table(s, 0.55, 1.65, 6.4, ["Radius", "Places", "Food", "Top categories"], rows,
           [1.5, 0.8, 0.8, 3.3], size=10)

    _text(s, 7.35, 1.65, 5.4, 0.4, "SCORE BREAKDOWN", 11, MUTED, bold=True)
    y = 2.12
    for c in r["score_components"]:
        _text(s, 7.35, y, 4.0, 0.3, c["label"], 12, WHITE)
        _text(s, 11.9, y, 0.9, 0.3, f"{c['score']}", 12, MINT, bold=True, align=PP_ALIGN.RIGHT)
        _bar(s, 7.35, y + 0.32, 5.45, 0.11, c["score"])
        y += 0.68

    b = r["block"]
    _text(s, 0.55, 5.55, 12.2, 0.9,
          f"{b['retail_continuity'][0]}. {b['retail_continuity'][1]}", 13, WHITE)
    _note(s, "Counts reflect places mapped in OpenStreetMap. Absence from the map is not evidence "
             "of absence.")


def _block(prs, r):
    b = r["block"]
    s = _slide(prs, "The block", b["street_axis_desc"])
    cols = [("One side", b["immediate_left"]), ("Other side", b["immediate_right"]),
            ("Across the street", b["across_street"])]
    x = 0.55
    for title, items in cols:
        _rect(s, x, 1.65, 3.95, 3.6, PANEL)
        _text(s, x + 0.22, 1.8, 3.5, 0.3, title.upper(), 10.5, MINT, bold=True)
        lines = [(n["name"][:26], f"{n['distance_ft']} ft") for n in items[:6]] or ["Nothing mapped"]
        _bullets(s, x + 0.2, 2.2, 3.6, 3.0, lines, size=12, spacing=6)
        x += 4.15

    facts = [
        ("Position:", "Corner" if b.get("is_corner") else "Mid-block"),
        ("Storefronts within ~240 ft:", str(b.get("storefront_count", 0))),
        ("Food businesses adjacent:", str(b.get("food_neighbours", 0))),
        ("Tagged vacancies:", str(b["vacancy"]["detected"])),
        ("Construction permits on street:", str(b.get("construction_count", 0))),
    ]
    _bullets(s, 0.58, 5.45, 12.2, 1.2, facts, size=13, spacing=3)
    _note(s, "Left/right/across is inferred from map geometry, not a physical survey.")


def _landscape(prs, r):
    c = r["competition"]
    s = _slide(prs, "Competitor landscape",
               f"{c['saturation_label']} - pressure index {c['pressure']}")
    _kpi(s, 0.55, 1.62, 2.9, 1.2, "Total competitors", str(c["total"]), WHITE, "within 0.5 mi")
    _kpi(s, 3.65, 1.62, 2.9, 1.2, "Direct competitors", str(c["direct_count"]), WHITE, "same category")
    _kpi(s, 6.75, 1.62, 2.9, 1.2, "Within one block", str(c["within_block"]), WHITE, "0.15 mi")
    _kpi(s, 9.85, 1.62, 2.9, 1.2, "Open 24 hours", str(c["open_24h_count"]), WHITE,
         "overnight coverage")

    rows = []
    for x in r["competitors"][:9]:
        rows.append([
            x["name"][:30], x["category_label"],
            f"{x['distance_mi']:.2f} mi" if x.get("distance_mi") else "-",
            f"{x['rating']}★" if x.get("rating") else "-",
            f"{x['reviews']:,}" if x.get("reviews") else "-",
            (x["threat"], THREAT_RGB.get(x["threat"], WHITE)),
        ])
    _table(s, 0.55, 3.1, 12.2, ["Business", "Type", "Distance", "Rating", "Reviews", "Threat"],
           rows, [3.6, 2.4, 1.4, 1.1, 1.3, 1.6], size=11)
    if not r["google_enabled"]:
        _note(s, "No Google Maps key connected - ratings, review counts and hours are unavailable.")


def _competitors(prs, r):
    top = [c for c in r["competitors"] if c["threat"] in ("VERY HIGH", "HIGH")][:3] \
        or r["competitors"][:3]
    s = _slide(prs, "Strongest competitors", "The operators most likely to take your customers")
    if not top:
        _text(s, 0.58, 2.2, 11, 1, "No significant competitors identified nearby.", 16, WHITE)
        return
    x = 0.55
    width = 12.2 / len(top) - 0.2
    for c in top:
        colour = THREAT_RGB.get(c["threat"], MUTED)
        _rect(s, x, 1.62, width, 4.8, PANEL)
        _rect(s, x, 1.62, width, 0.1, colour)
        _text(s, x + 0.22, 1.86, width - 0.4, 0.6, c["name"][:34], 17, WHITE, bold=True)
        _text(s, x + 0.22, 2.5, width - 0.4, 0.32, c["threat"] + " THREAT", 12, colour, bold=True)
        facts = []
        if c.get("distance_mi") is not None:
            facts.append(f"{c['distance_mi']:.2f} mi - {c.get('walk_minutes', '?')} min walk")
        if c.get("rating"):
            facts.append(f"{c['rating']}★  {c.get('reviews', 0):,} reviews")
        if c.get("open_24h"):
            facts.append("Open 24 hours")
        if c.get("since"):
            facts.append(c["since"]["statement"])
        facts.append(c["category_label"])
        _bullets(s, x + 0.22, 2.95, width - 0.4, 1.3, facts, size=12, colour=MUTED, spacing=4)
        why = c["why"]
        _text(s, x + 0.22, 4.3, width - 0.4, 2.0, why[:290], 11.5, WHITE)
        x += width + 0.2
    _note(s, "Threat levels are modeled from distance, review volume, rating and operating hours.")


def _generators(prs, r):
    g = r["generators"]
    s = _slide(prs, "Demand generators", "What puts people on this sidewalk")
    rows = [[x["name"][:34], x["category_label"], f"{x['distance_ft']} ft",
             f"{x['walk_minutes']} min", f"{x['relevance']}/100"] for x in g["items"][:9]]
    _table(s, 0.55, 1.65, 7.3, ["Generator", "Type", "Distance", "Walk", "Relevance"], rows,
           [2.9, 1.7, 1.0, 0.85, 0.95], size=10.5)

    _text(s, 8.15, 1.65, 4.6, 0.4, "CATEGORY STRENGTH", 11, MUTED, bold=True)
    y = 2.1
    top_cats = sorted(g["strength"].items(), key=lambda kv: -kv[1])[:7]
    for cat, val in top_cats:
        _text(s, 8.15, y, 3.4, 0.3, cat.replace("_", " ").title(), 12, WHITE)
        _text(s, 11.95, y, 0.8, 0.3, str(val), 12, MINT, bold=True, align=PP_ALIGN.RIGHT)
        _bar(s, 8.15, y + 0.31, 4.6, 0.1, val)
        y += 0.62
    if g["top"]:
        _text(s, 0.55, 6.3, 7.3, 0.6, g["top"][0]["why"][:190], 11, MUTED)


def _customers(prs, r):
    c = r["customers"]
    s = _slide(prs, "Customer profile", "Who actually walks in, ranked")
    y = 1.7
    for x in c["ranked"][:5]:
        _rect(s, 0.55, y, 12.2, 0.92, PANEL)
        _text(s, 0.78, y + 0.1, 4.0, 0.35, x["label"], 15, WHITE, bold=True)
        _text(s, 0.78, y + 0.46, 4.0, 0.3, x["tier"], 11, MINT, bold=True)
        _text(s, 4.9, y + 0.13, 7.6, 0.7, x["buys"][:180], 11, MUTED)
        _bar(s, 0.78, y + 0.78, 3.6, 0.08, x["score"])
        y += 1.02
    _text(s, 0.55, 6.5, 12.2, 0.5,
          f"Customer mix: {c['concentration']['level']}. {c['concentration']['note']}", 11.5, WHITE)


def _dayparts(prs, r):
    d = r["dayparts"]
    s = _slide(prs, "Dayparts", f"MODELED demand by time of day - {d['shape'][0]}")
    y = 1.68
    for w in d["windows"]:
        _text(s, 0.55, y, 1.8, 0.3, w["label"], 12.5, WHITE, bold=True)
        _bar(s, 2.5, y + 0.06, 4.2, 0.18, w["weekday"], MINT)
        _text(s, 6.8, y, 0.7, 0.3, str(w["weekday"]), 11, MUTED)
        _bar(s, 7.6, y + 0.06, 4.2, 0.18, w["weekend"], RGBColor(0x7E, 0xA8, 0xFF))
        _text(s, 11.9, y, 0.8, 0.3, str(w["weekend"]), 11, MUTED)
        y += 0.6
    _text(s, 2.5, 1.35, 4.2, 0.3, "WEEKDAY", 10, MINT, bold=True)
    _text(s, 7.6, 1.35, 4.2, 0.3, "WEEKEND", 10, RGBColor(0x7E, 0xA8, 0xFF), bold=True)
    _text(s, 0.55, 6.5, 12.2, 0.45,
          f"Peak: {d['peak']['label']}. Weakest: {d['trough']['label']}. {d['shape'][1]}", 12, WHITE)
    _note(s, d["disclaimer"][:170])


def _demographics(prs, r):
    d = r["demographics"]
    s = _slide(prs, "Demographics", f"Census tract {d.get('tract_name') or 'unavailable'}")
    if not d.get("available"):
        _text(s, 0.58, 2.2, 11, 1, "Census data unavailable for this location. UNKNOWN, not zero.",
              18, WHITE)
        return
    f = d["facts"]

    def val(key, fmt):
        fv = f.get(key) or {}
        return fmt.format(fv["value"]) if fv.get("known") else "N/A"

    _kpi(s, 0.55, 1.65, 3.9, 1.35, "Tract population", val("population", "{:,.0f}"), WHITE)
    _kpi(s, 4.7, 1.65, 3.9, 1.35, "Median household income", val("median_income", "${:,.0f}"), MINT)
    _kpi(s, 8.85, 1.65, 3.9, 1.35, "Population density", val("density_sq_mi", "{:,.0f}"), WHITE,
         "per square mile")
    _kpi(s, 0.55, 3.2, 3.9, 1.35, "Households", val("households", "{:,.0f}"), WHITE)
    _kpi(s, 4.7, 3.2, 3.9, 1.35, "Median age", val("median_age", "{:.1f}"), WHITE)
    _kpi(s, 8.85, 3.2, 3.9, 1.35, "Average household size", val("avg_household_size", "{:.2f}"), WHITE)

    rows = []
    for radius, fv in sorted((d.get("radius_estimates") or {}).items()):
        rows.append([f"{radius} mi radius",
                     f"{fv['value']:,}" if fv.get("known") else "N/A", "INFERRED from density"])
    if rows:
        _table(s, 0.55, 4.8, 7.0, ["Geography", "Estimated residents", "Evidence"], rows,
               [2.2, 2.2, 2.6], size=11)
    ch = (d.get("character") or {}).get("value")
    if ch:
        _text(s, 8.0, 4.9, 4.7, 1.2, ch, 14, WHITE)
    _note(s, "Tract totals are tract totals. Radius figures are derived estimates, not counts.")


def _sales(prs, r):
    sm = r["sales"]
    s = _slide(prs, "Sales scenarios", "MODELED estimate built from customer pools")
    order = ["Conservative", "Realistic", "Strong Operator", "Elite Operator"]
    rows = []
    for name in order:
        x = sm["scenarios"][name]
        rows.append([name, f"${x['daily']:,}", f"${x['weekly']:,}", f"${x['monthly']:,}",
                     f"${x['annual']:,}"])
    _table(s, 0.55, 1.65, 7.2, ["Scenario", "Daily", "Weekly", "Monthly", "Annual"], rows,
           [1.9, 1.2, 1.3, 1.4, 1.4], size=12)

    _text(s, 8.0, 1.65, 4.8, 0.35, "PUSHING IT UP", 11, MINT, bold=True)
    ups = [(x["factor"], x["effect"]) for x in sm["drivers_up"][:4]] or ["None identified"]
    _bullets(s, 8.0, 2.05, 4.8, 1.7, ups, size=12, spacing=5)
    _text(s, 8.0, 3.9, 4.8, 0.35, "PULLING IT DOWN", 11, SEV_RGB["HIGH"], bold=True)
    downs = [(x["factor"], x["effect"]) for x in sm["drivers_down"][:4]] or ["None identified"]
    _bullets(s, 8.0, 4.3, 4.8, 1.7, downs, size=12, spacing=5)

    lines = [(l["label"], f"{l['pool_size']:,} people, {l['effective_pct']}% capture, "
                          f"{l['transactions']:.0f} sales/day") for l in sm["transaction_lines"][:5]]
    _text(s, 0.55, 3.9, 7.2, 0.35, "HOW THE NUMBER IS BUILT", 11, MUTED, bold=True)
    _bullets(s, 0.55, 4.28, 7.2, 2.0, lines, size=11.5, spacing=5)
    _note(s, "ESTIMATE. Verify against the seller's register tapes and sales tax filings.")


def _economics(prs, r):
    pnl = r["pnl"]["by_scenario"]
    ra = r["rent_assessment"]
    s = _slide(prs, "Rent and economics",
               f"Rent {ra['rent_pct']}% of modeled sales - {ra['band']}")
    order = ["Conservative", "Realistic", "Strong Operator"]
    fields = [("revenue", "Sales"), ("cogs", "Cost of goods"), ("labour", "Labour"),
              ("delivery_commission", "Delivery commission"), ("rent", "Rent"),
              ("fixed", "Other fixed"), ("profit", "OPERATING PROFIT")]
    rows = []
    for key, label in fields:
        row = [label]
        for o in order:
            v = pnl[o][key]
            colour = WHITE
            if key == "profit":
                colour = MINT if v > 0 else SEV_RGB["CRITICAL"]
            row.append((f"${v:,}", colour))
        rows.append(row)
    _table(s, 0.55, 1.65, 8.0, ["Monthly"] + order, rows, [2.6, 1.8, 1.8, 1.8], size=12)

    _kpi(s, 8.9, 1.65, 3.85, 1.3, "Break-even sales",
         f"${r['pnl']['breakeven_daily']:,}/day" if r["pnl"].get("breakeven_daily") else "n/a",
         WHITE)
    _kpi(s, 8.9, 3.1, 3.85, 1.3, "Max supportable rent",
         f"${ra['max_supportable_rent']:,}", MINT, "per month at modeled sales")
    _kpi(s, 8.9, 4.55, 3.85, 1.3, "Annual profit (realistic)",
         f"${pnl['Realistic']['annual_profit']:,}",
         MINT if pnl["Realistic"]["annual_profit"] > 0 else SEV_RGB["CRITICAL"])
    _text(s, 0.55, 5.7, 8.0, 0.8, ra["verdict"], 12, WHITE)
    _note(s, "Excludes debt service, buildout amortisation, owner salary and taxes."
             + (" Rent is ESTIMATED." if r["rent"]["estimated"] else ""))


def _opportunity(prs, r):
    o = r["opportunities"]
    s = _slide(prs, "The opportunity", "Gaps confirmed by more than one signal")
    if not o["items"]:
        _text(s, 0.58, 2.2, 11.5, 1.5,
              "No multi-signal opportunity was detected. This location would be won on execution "
              "alone, not on an unserved gap in the market.", 18, WHITE)
        return
    y = 1.68
    for x in o["items"][:4]:
        colour = MINT if x["strength"] == "HIGH" else RGBColor(0xE0, 0xA0, 0x2A)
        _rect(s, 0.55, y, 12.2, 1.18, PANEL)
        _rect(s, 0.55, y, 0.08, 1.18, colour)
        _text(s, 0.82, y + 0.1, 9.5, 0.4, x["title"], 15, WHITE, bold=True)
        _text(s, 11.4, y + 0.12, 1.3, 0.3, x["strength"], 11, colour, bold=True)
        _text(s, 0.82, y + 0.52, 11.6, 0.6, x["action"][:210], 11.5, MUTED)
        y += 1.28


def _risks(prs, r):
    risks = r["risks"]
    s = _slide(prs, "Risks", "What could kill this deal")
    if not risks["items"]:
        _text(s, 0.58, 2.2, 11, 1, "No material red flags detected.", 18, WHITE)
        return
    y = 1.68
    for x in risks["items"][:5]:
        colour = SEV_RGB.get(x["severity"], MUTED)
        _rect(s, 0.55, y, 12.2, 0.95, PANEL)
        _rect(s, 0.55, y, 0.08, 0.95, colour)
        _text(s, 0.82, y + 0.08, 9.4, 0.35, x["title"], 14.5, WHITE, bold=True)
        _text(s, 11.3, y + 0.1, 1.4, 0.3, x["severity"], 11, colour, bold=True)
        _text(s, 0.82, y + 0.46, 11.6, 0.45, x["detail"][:180], 11, MUTED)
        y += 1.05
    counts = " | ".join(f"{k}: {v}" for k, v in risks["counts"].items())
    _note(s, counts)


def _final(prs, r, photos):
    s = _slide(prs, "Final recommendation")
    colour = VERDICT_RGB.get(r["verdict"], MINT)
    _rect(s, 0.55, 1.5, 12.2, 1.15, colour)
    _text(s, 0.85, 1.68, 11.6, 0.8, r["five_hundred_k"]["headline"], 22,
          RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    _text(s, 0.58, 2.85, 7.4, 2.2, r["five_hundred_k"]["explanation"], 13, WHITE)

    sv = photos.get("streetview")
    if sv:
        try:
            s.shapes.add_picture(io.BytesIO(sv), Inches(8.3), Inches(2.85),
                                 width=Inches(4.45), height=Inches(2.78))
            _text(s, 8.3, 5.7, 4.45, 0.3,
                  f"Street View, {(r.get('streetview') or {}).get('date', 'date unknown')}",
                  9, MUTED, italic=True)
        except Exception:  # noqa: BLE001
            pass

    _text(s, 0.58, 5.3, 7.4, 0.35, "BEFORE YOU SIGN", 11, MINT, bold=True)
    _bullets(s, 0.58, 5.65, 7.4, 1.2, r["checklist"][:3], size=11, colour=MUTED, spacing=3)
    _note(s, f"Data confidence {r['confidence']['score']}/100. This deck is a screening tool, "
             "not a substitute for walking the block.")


# ------------------------------------------------------------- comparison
def build_comparison(comparison):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    s = _slide(prs)
    _rect(s, 0, 0, Inches(0.13), H, MINT)
    _text(s, 0.75, 2.0, 11.5, 1.0, "Location comparison", 40, WHITE, bold=True)
    _text(s, 0.78, 3.2, 11.5, 0.5,
          f"{comparison['count']} locations ranked - {comparison.get('concept_label', '')}",
          16, MUTED)
    if comparison.get("winner"):
        _text(s, 0.78, 4.2, 11.5, 0.8, comparison["winner"]["headline"], 20, MINT, bold=True)

    s = _slide(prs, "Ranking", "Sorted by overall score")
    rows = []
    for r in comparison["rows"][:12]:
        rows.append([f"#{r['rank']}", r["address"][:34], (str(r["score"]), MINT),
                     str(r["confidence"]), f"${r['realistic_daily']:,}",
                     f"${r['monthly_profit']:,}", f"{r['rent_pct']}%",
                     (r["verdict"], VERDICT_RGB.get(r["verdict"], WHITE))])
    _table(s, 0.4, 1.6, 12.5, ["#", "Address", "Score", "Conf", "Sales/day", "Profit/mo",
                               "Rent %", "Verdict"], rows,
           [0.5, 3.4, 0.8, 0.7, 1.4, 1.4, 0.9, 1.5], size=10.5)

    if comparison.get("winner"):
        w = comparison["winner"]
        s = _slide(prs, "The one to take")
        _rect(s, 0.55, 1.5, 12.2, 1.1, MINT)
        _text(s, 0.85, 1.68, 11.6, 0.8, w["headline"], 22, RGBColor(0x06, 0x22, 0x18), bold=True)
        _text(s, 0.58, 2.85, 12.2, 2.4, w["why"], 14, WHITE)
        _text(s, 0.58, 5.6, 12.2, 0.8, w["caveat"], 12, MUTED, italic=True)

    for r in comparison["rows"][:8]:
        s = _slide(prs, r["address"], f"Rank #{r['rank']} - {r['verdict']}")
        _kpi(s, 0.55, 1.7, 2.9, 1.3, "Score", str(r["score"]),
             VERDICT_RGB.get(r["verdict"], WHITE))
        _kpi(s, 3.65, 1.7, 2.9, 1.3, "Realistic sales", f"${r['realistic_daily']:,}", MINT, "per day")
        _kpi(s, 6.75, 1.7, 2.9, 1.3, "Monthly profit", f"${r['monthly_profit']:,}",
             MINT if r["monthly_profit"] > 0 else SEV_RGB["CRITICAL"])
        _kpi(s, 9.85, 1.7, 2.9, 1.3, "Rent", f"${r['rent']:,}", WHITE, f"{r['rent_pct']}% of sales")
        _bullets(s, 0.58, 3.35, 12.2, 2.6, [
            ("Competition:", f"{r['direct_competitors']} direct of {r['competitors']} total - "
                             f"{r['saturation']}"),
            ("Strongest competitor:", r.get("strongest_competitor") or "None identified"),
            ("Biggest opportunity:", r["opportunity"]),
            ("Biggest risk:", f"{r['risk']} ({r['risk_severity']})"),
            ("Data confidence:", f"{r['confidence']}/100"),
        ], size=14, spacing=10)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
